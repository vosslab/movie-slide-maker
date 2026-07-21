"""Convert movie presentations to validated OpenDocument presentations."""

# Standard Library
import os
import re
import pathlib
import zipfile
import tempfile
import subprocess
import xml.etree.ElementTree

# PIP3 modules
import pptx

# local repo modules
import slide_maker.slide_builder


DRAW_NS = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
FO_NS = "urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0"
OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
PRESENTATION_NS = "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0"
STYLE_NS = "urn:oasis:names:tc:opendocument:xmlns:style:1.0"
SVG_NS = "urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0"
XLINK_NS = "http://www.w3.org/1999/xlink"
REQUIRED_ROLES = (
	slide_maker.slide_builder.TEMPLATE_TITLE_NAME,
	slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME,
	slide_maker.slide_builder.TEMPLATE_POSTER_NAME,
)
REQUIRED_FIELDS = (
	"IMDB rating",
	"Critics: RT",
	"Genre:",
	"Director:",
	"Run time:",
	"Review Summary:",
)
GEOMETRY_TOLERANCE_CM = 0.1


class SlideConversionError(RuntimeError):
	"""Report a failed conversion or invalid converted presentation."""


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise a conversion error when a required condition is false."""
	if condition:
		return
	raise SlideConversionError(message)


#============================================
def parse_xml(archive: zipfile.ZipFile, member_name: str) -> xml.etree.ElementTree.Element:
	"""Parse one required XML member from an ODP package."""
	require(member_name in archive.namelist(), f"Converted ODP is missing {member_name}")
	# LibreOffice creates this package locally from the builder's scratch presentation.
	root = xml.etree.ElementTree.fromstring(archive.read(member_name))  # nosec B314
	return root


#============================================
def length_cm(value: str) -> float:
	"""Convert an ODF length with a supported physical unit to centimeters."""
	match = re.fullmatch(r"([+-]?[0-9]+(?:\.[0-9]+)?)(cm|mm|in|pt)", value)
	require(match is not None, f"Converted ODP has an unsupported length: {value}")
	amount = float(match.group(1))
	unit = match.group(2)
	factors = {"cm": 1.0, "mm": 0.1, "in": 2.54, "pt": 2.54 / 72.0}
	centimeters = amount * factors[unit]
	return centimeters


#============================================
def named_frames(page: xml.etree.ElementTree.Element) -> dict[str, xml.etree.ElementTree.Element]:
	"""Return movie frames indexed by their semantic draw names."""
	name_attribute = f"{{{DRAW_NS}}}name"
	frames: dict[str, xml.etree.ElementTree.Element] = {}
	for frame in page.iter(f"{{{DRAW_NS}}}frame"):
		name = frame.get(name_attribute)
		if name in REQUIRED_ROLES:
			require(name not in frames, f"Converted movie slide has ambiguous role: {name}")
			frames[name] = frame
	return frames


#============================================
def page_is_hidden(
	page: xml.etree.ElementTree.Element,
	content_root: xml.etree.ElementTree.Element,
) -> bool:
	"""Return whether an ODP page or its drawing-page style is hidden."""
	visibility_attribute = f"{{{PRESENTATION_NS}}}visibility"
	if page.get(visibility_attribute) == "hidden":
		return True
	style_name = page.get(f"{{{DRAW_NS}}}style-name")
	if not style_name:
		return False
	for style in content_root.iter(f"{{{STYLE_NS}}}style"):
		if style.get(f"{{{STYLE_NS}}}name") != style_name:
			continue
		properties = style.find(f"{{{STYLE_NS}}}drawing-page-properties")
		if properties is not None:
			return properties.get(visibility_attribute) == "hidden"
	return False


#============================================
def presentation_pages(
	content_root: xml.etree.ElementTree.Element,
) -> list[xml.etree.ElementTree.Element]:
	"""Return only top-level presentation pages, excluding nested notes pages."""
	presentations = list(content_root.iter(f"{{{OFFICE_NS}}}presentation"))
	require(len(presentations) == 1, "Converted ODP has no single presentation body")
	pages = presentations[0].findall(f"{{{DRAW_NS}}}page")
	require(bool(pages), "Converted ODP presentation body has no slides")
	return pages


#============================================
def movie_pages(
	content_root: xml.etree.ElementTree.Element,
	expected_movie_slides: int = 1,
) -> list[xml.etree.ElementTree.Element]:
	"""Return all and only the expected visible movie slides in document order."""
	require(expected_movie_slides > 0, "Expected movie-slide count must be positive")
	pages = []
	hidden_pages = []
	for page in presentation_pages(content_root):
		if page_is_hidden(page, content_root):
			hidden_pages.append(page)
			continue
		roles = set(named_frames(page))
		require(
			roles == set(REQUIRED_ROLES),
			f"Converted ODP has a visible non-movie slide with roles: {sorted(roles)!r}",
		)
		pages.append(page)
	require(
		len(hidden_pages) == 1,
		f"Converted ODP has {len(hidden_pages)} hidden slides; expected one template",
	)
	require(
		len(pages) == expected_movie_slides,
		f"Converted ODP has {len(pages)} visible movie slides; "
		f"expected {expected_movie_slides}",
	)
	return pages


#============================================
def validate_landscape(
	page: xml.etree.ElementTree.Element,
	styles_root: xml.etree.ElementTree.Element,
) -> None:
	"""Validate the converted movie slide's referenced page layout is landscape."""
	master_name = page.get(f"{{{DRAW_NS}}}master-page-name")
	require(bool(master_name), "Converted movie slide has no master-page reference")
	master_pages = [
		master for master in styles_root.iter(f"{{{STYLE_NS}}}master-page")
		if master.get(f"{{{STYLE_NS}}}name") == master_name
	]
	require(len(master_pages) == 1, "Converted movie slide master-page reference is invalid")
	layout_name = master_pages[0].get(f"{{{STYLE_NS}}}page-layout-name")
	require(bool(layout_name), "Converted movie slide master has no page layout")
	layouts = [
		layout for layout in styles_root.iter(f"{{{STYLE_NS}}}page-layout")
		if layout.get(f"{{{STYLE_NS}}}name") == layout_name
	]
	require(len(layouts) == 1, "Converted movie slide page-layout reference is invalid")
	properties = layouts[0].find(f"{{{STYLE_NS}}}page-layout-properties")
	require(properties is not None, "Converted movie slide page layout has no properties")
	width_value = properties.get(f"{{{FO_NS}}}page-width")
	height_value = properties.get(f"{{{FO_NS}}}page-height")
	require(bool(width_value and height_value), "Converted movie slide has no page dimensions")
	width_cm = length_cm(width_value)
	height_cm = length_cm(height_value)
	require(width_cm > height_cm, "Converted movie slide is not landscape")


#============================================
def validate_source_displayed_fields(source_shapes: dict[str, object]) -> None:
	"""Validate required title, labels, and poster payload in one source slide."""
	title = source_shapes[slide_maker.slide_builder.TEMPLATE_TITLE_NAME]
	outline = source_shapes[slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME]
	poster = source_shapes[slide_maker.slide_builder.TEMPLATE_POSTER_NAME]
	require(bool(title.text.strip()), "Source movie slide title is empty")
	for field in REQUIRED_FIELDS:
		require(field in outline.text, f"Source movie slide is missing displayed field: {field}")
	require(bool(poster.image.blob), "Source movie slide poster payload is empty")


#============================================
def source_movie_shape_sets(
	pptx_path: pathlib.Path,
	expected_movie_slides: int = 1,
) -> list[dict[str, object]]:
	"""Return the expected source movie-shape sets in slide order."""
	require(expected_movie_slides > 0, "Expected movie-slide count must be positive")
	presentation = pptx.Presentation(pptx_path)
	require(
		presentation.slide_width > presentation.slide_height,
		"Source movie presentation is not landscape",
	)
	shape_sets = []
	hidden_slides = []
	for slide in presentation.slides:
		if slide._element.get("show") == "0":
			hidden_slides.append(slide)
			continue
		shapes = {shape.name: shape for shape in slide.shapes if shape.name in REQUIRED_ROLES}
		require(
			set(shapes) == set(REQUIRED_ROLES),
			f"Source PPTX has a visible non-movie slide with roles: {sorted(shapes)!r}",
		)
		validate_source_displayed_fields(shapes)
		shape_sets.append(shapes)
	require(
		len(hidden_slides) == 1,
		f"Source PPTX has {len(hidden_slides)} hidden slides; expected one template",
	)
	template_roles = {
		shape.name for shape in hidden_slides[0].shapes if shape.name in REQUIRED_ROLES
	}
	require(
		template_roles == set(REQUIRED_ROLES),
		f"Source PPTX hidden template has invalid roles: {sorted(template_roles)!r}",
	)
	require(
		len(shape_sets) == expected_movie_slides,
		f"Source PPTX has {len(shape_sets)} visible movie slides; "
		f"expected {expected_movie_slides}",
	)
	return shape_sets


#============================================
def validate_geometry(
	frames: dict[str, xml.etree.ElementTree.Element],
	source_shapes: dict[str, object],
) -> None:
	"""Compare converted frame geometry with the source using the plan tolerance."""
	attributes = (("x", "left"), ("y", "top"), ("width", "width"), ("height", "height"))
	for role in REQUIRED_ROLES:
		for odf_name, source_name in attributes:
			value = frames[role].get(f"{{{SVG_NS}}}{odf_name}")
			require(bool(value), f"Converted role {role} has no {odf_name} geometry")
			converted_cm = length_cm(value)
			source_cm = float(getattr(source_shapes[role], source_name)) / 360000.0
			difference = abs(converted_cm - source_cm)
			require(
				difference <= GEOMETRY_TOLERANCE_CM,
				f"Converted role {role} {odf_name} differs by {difference:.3f} cm",
			)


#============================================
def validate_displayed_fields(page: xml.etree.ElementTree.Element) -> None:
	"""Validate required title, labels, and poster payload on the movie slide."""
	frames = named_frames(page)
	title_text = "".join(
		frames[slide_maker.slide_builder.TEMPLATE_TITLE_NAME].itertext()
	).strip()
	require(bool(title_text), "Converted movie slide title is empty")
	outline_text = " ".join(
		frames[slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME].itertext()
	)
	for field in REQUIRED_FIELDS:
		require(field in outline_text, f"Converted movie slide is missing displayed field: {field}")


#============================================
def normalized_frame_text(frame: xml.etree.ElementTree.Element) -> str:
	"""Return whitespace-normalized text from one converted semantic frame."""
	plain_text = " ".join("".join(frame.itertext()).split())
	return plain_text


#============================================
def validate_poster(
	archive: zipfile.ZipFile,
	page: xml.etree.ElementTree.Element,
) -> None:
	"""Validate that the poster role references a nonempty packaged image."""
	poster = named_frames(page)[slide_maker.slide_builder.TEMPLATE_POSTER_NAME]
	images = list(poster.iter(f"{{{DRAW_NS}}}image"))
	require(len(images) == 1, "Converted movie slide poster role has no single image")
	href = images[0].get(f"{{{XLINK_NS}}}href")
	require(bool(href), "Converted movie slide poster has no package reference")
	member_name = href.removeprefix("./")
	require(member_name in archive.namelist(), "Converted movie slide poster payload is absent")
	require(archive.getinfo(member_name).file_size > 0, "Converted movie slide poster is empty")


#============================================
def validate_odp(
	odp_path: pathlib.Path,
	source_pptx_path: pathlib.Path,
	expected_movie_slides: int = 1,
) -> None:
	"""Validate every paired source and converted movie slide in order."""
	require(odp_path.is_file(), f"LibreOffice did not create the expected ODP: {odp_path}")
	require(zipfile.is_zipfile(odp_path), f"Converted output is not an ODP package: {odp_path}")
	source_shape_sets = source_movie_shape_sets(source_pptx_path, expected_movie_slides)
	with zipfile.ZipFile(odp_path) as archive:
		content_root = parse_xml(archive, "content.xml")
		styles_root = parse_xml(archive, "styles.xml")
		pages = movie_pages(content_root, expected_movie_slides)
		for index, (page, source_shapes) in enumerate(zip(pages, source_shape_sets, strict=True), 1):
			frames = named_frames(page)
			validate_landscape(page, styles_root)
			validate_displayed_fields(page)
			validate_poster(archive, page)
			validate_geometry(frames, source_shapes)
			converted_title = normalized_frame_text(
				frames[slide_maker.slide_builder.TEMPLATE_TITLE_NAME]
			)
			source_title = " ".join(
				source_shapes[slide_maker.slide_builder.TEMPLATE_TITLE_NAME].text.split()
			)
			require(
				converted_title == source_title,
				f"Converted movie slide {index} changed title identity: {converted_title}",
			)


#============================================
def convert_presentation(
	scratch_pptx_path: pathlib.Path,
	output_odp_path: pathlib.Path,
	expected_movie_slides: int = 1,
) -> pathlib.Path:
	"""Convert and validate expected movie slides, then remove accepted scratch."""
	require(scratch_pptx_path.is_file(), f"Scratch PPTX is absent: {scratch_pptx_path}")
	output_odp_path.parent.mkdir(parents=True, exist_ok=True)
	with tempfile.TemporaryDirectory() as temporary_directory:
		conversion_dir = pathlib.Path(temporary_directory)
		profile_uri = conversion_dir.joinpath("libreoffice_profile").resolve().as_uri()
		command = [
			"soffice",
			f"-env:UserInstallation={profile_uri}",
			"--headless",
			"--convert-to",
			"odp",
			"--outdir",
			str(conversion_dir),
			str(scratch_pptx_path.resolve()),
		]
		result = subprocess.run(command, capture_output=True, text=True, check=False)
		converted_path = conversion_dir / f"{scratch_pptx_path.stem}.odp"
		diagnostic = result.stderr.strip() or result.stdout.strip() or "no LibreOffice diagnostic"
		require(result.returncode == 0, f"LibreOffice conversion failed: {diagnostic}")
		validate_odp(converted_path, scratch_pptx_path, expected_movie_slides)
		os.replace(converted_path, output_odp_path)
	scratch_pptx_path.unlink()
	return output_odp_path
