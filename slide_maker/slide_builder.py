"""Build movie slides from the extracted presentation template."""

# Standard Library
import copy
import pathlib

# PIP3 modules
import pptx
import PIL.Image #pillow
import pptx.util
import pptx.oxml
import pptx.oxml.ns
import pptx.enum.text

# local repo modules
import slide_maker.moviedata
import slide_maker.emoji_marks


TEMPLATE_TITLE_NAME = "movie_title"
TEMPLATE_OUTLINE_NAME = "movie_outline"
TEMPLATE_POSTER_NAME = "movie_poster"
FONT_NAME = "OpenDyslexic"
TITLE_FONT_SIZE = pptx.util.Pt(40.4)
PRIMARY_FONT_SIZE = pptx.util.Pt(22.0)
SECONDARY_FONT_SIZE = pptx.util.Pt(19.2)
PAGE_WIDTH = pptx.util.Cm(28.0)
PAGE_HEIGHT = pptx.util.Cm(17.5)
IMDB_SCORE_HIGHLIGHT = "F5C518"


class SlideBuildError(RuntimeError):
	"""Report a missing or ambiguous template role required for a movie slide."""


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise a builder error when a required template condition is false."""
	if condition:
		return
	raise SlideBuildError(message)


#============================================
def shape_by_name(slide: object, name: str) -> object:
	"""Return the one shape carrying a required semantic role name."""
	candidates = [shape for shape in slide.shapes if shape.name == name]
	require(candidates, f"Required template role is absent: {name}")
	require(len(candidates) == 1, f"Required template role is ambiguous: {name}")
	shape = candidates[0]
	return shape


#============================================
def clone_template_shapes(source_slide: object, target_slide: object) -> None:
	"""Copy template shapes to a new slide while retaining their formatting."""
	for shape in list(target_slide.shapes):
		target_slide.shapes._spTree.remove(shape._element)
	for shape in source_slide.shapes:
		cloned_element = copy.deepcopy(shape._element)
		target_slide.shapes._spTree.insert_element_before(cloned_element, "p:extLst")


#============================================
def format_compact_count(count: int) -> str:
	"""Format a nonnegative count with a compact suffix and at least two significant digits."""
	if type(count) is not int or count < 0:
		raise ValueError("Compact count must be a nonnegative integer")
	units = (
		(999_500_000, 1_000_000_000, "B"),
		(999_500, 1_000_000, "M"),
		(1_000, 1_000, "k"),
	)
	for threshold, divisor, suffix in units:
		if count < threshold:
			continue
		scaled = count / divisor
		if scaled < 10:
			number = f"{scaled:.1f}"
		else:
			number = f"{scaled:.0f}"
		return f"{number}{suffix}"
	return str(count)


#============================================
def add_styled_run(
	paragraph: object,
	text: str,
	font_size: pptx.util.Length,
	bold: bool = False,
) -> object:
	"""Add one consistently styled OpenDyslexic run."""
	run = paragraph.add_run()
	run.text = text
	run.font.name = FONT_NAME
	run.font.size = font_size
	run.font.bold = bold
	return run


#============================================
def set_run_highlight(run: object, color: str) -> None:
	"""Apply one DrawingML text highlight not exposed by python-pptx."""
	run_properties = run._r.get_or_add_rPr()
	highlight_tag = pptx.oxml.ns.qn("a:highlight")
	for child in list(run_properties):
		if child.tag == highlight_tag:
			run_properties.remove(child)
	highlight_xml = (
		f'<a:highlight {pptx.oxml.ns.nsdecls("a")}>'
		f'<a:srgbClr val="{color}"/>'
		"</a:highlight>"
	)
	highlight = pptx.oxml.parse_xml(highlight_xml)
	run_properties.append(highlight)


#============================================
def add_text_paragraph(
	text_frame: object,
	text: str,
	level: int,
	font_size: pptx.util.Length,
	first: bool,
) -> None:
	"""Add one OpenDyslexic paragraph at the requested outline level."""
	if first:
		paragraph = text_frame.paragraphs[0]
	else:
		paragraph = text_frame.add_paragraph()
	paragraph.level = level
	add_styled_run(paragraph, text, font_size)


#============================================
def add_imdb_paragraph(
	text_frame: object,
	movie_data: slide_maker.moviedata.MovieData,
) -> None:
	"""Add the IMDb line with its score highlighted in brand yellow."""
	paragraph = text_frame.add_paragraph()
	paragraph.level = 1
	add_styled_run(paragraph, "IMDB rating ", SECONDARY_FONT_SIZE)
	score_run = add_styled_run(
		paragraph,
		f"{movie_data.imdb_rating:.1f}",
		SECONDARY_FONT_SIZE,
		True,
	)
	set_run_highlight(score_run, IMDB_SCORE_HIGHLIGHT)
	vote_text = format_compact_count(movie_data.imdb_votes)
	add_styled_run(paragraph, f", {vote_text} votes", SECONDARY_FONT_SIZE)


#============================================
def format_rating_marks(
	movie_data: slide_maker.moviedata.MovieData,
) -> tuple[str, str | None, str]:
	"""Return the requested critic, audience, and Metascore display marks."""
	rt_mark = slide_maker.emoji_marks.rt_critic_mark_for_score(movie_data.rt_tomatometer)
	if movie_data.rt_audience_score is None:
		audience_mark = None
	else:
		audience_mark = slide_maker.emoji_marks.rt_audience_mark_for_score(
			movie_data.rt_audience_score
		)
	metascore_marks = {
		"high": slide_maker.emoji_marks.GREEN_SQUARE_MARK,
		"middle": slide_maker.emoji_marks.YELLOW_SQUARE_MARK,
		"low": slide_maker.emoji_marks.RED_SQUARE_MARK,
	}
	metascore_mark = metascore_marks[movie_data.metascore_band]
	return rt_mark, audience_mark, metascore_mark


#============================================
def fill_title(slide: object, movie_data: slide_maker.moviedata.MovieData) -> None:
	"""Fill the named movie-title anchor with the product title line."""
	title_shape = shape_by_name(slide, TEMPLATE_TITLE_NAME)
	require(title_shape.has_text_frame, "Movie title role has no text frame")
	text_frame = title_shape.text_frame
	text_frame.clear()
	text_frame.word_wrap = True
	text_frame.vertical_anchor = pptx.enum.text.MSO_VERTICAL_ANCHOR.MIDDLE
	text_frame.auto_size = pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
	title_text = f"{movie_data.title} ({movie_data.year})"
	add_text_paragraph(text_frame, title_text, 0, TITLE_FONT_SIZE, True)


#============================================
def fill_outline(slide: object, movie_data: slide_maker.moviedata.MovieData) -> None:
	"""Fill the named outline anchor with product labels and bullet hierarchy."""
	outline_shape = shape_by_name(slide, TEMPLATE_OUTLINE_NAME)
	require(outline_shape.has_text_frame, "Movie outline role has no text frame")
	rt_mark, audience_mark, metascore_mark = format_rating_marks(movie_data)
	if movie_data.rt_audience_score is None:
		audience_text = "Audience: N/A"
	else:
		audience_text = f"Audience: {audience_mark} {movie_data.rt_audience_score}%"
	critics_text = (
		f"Critics: RT {rt_mark} {movie_data.rt_tomatometer}% / "
		f"MS {metascore_mark} {movie_data.metascore}"
	)
	text_frame = outline_shape.text_frame
	text_frame.clear()
	text_frame.word_wrap = True
	text_frame.auto_size = pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
	add_text_paragraph(text_frame, movie_data.plot, 0, PRIMARY_FONT_SIZE, True)
	add_imdb_paragraph(text_frame, movie_data)
	add_text_paragraph(text_frame, critics_text, 1, SECONDARY_FONT_SIZE, False)
	add_text_paragraph(text_frame, audience_text, 1, SECONDARY_FONT_SIZE, False)
	add_text_paragraph(
		text_frame,
		f"Genre: {', '.join(movie_data.genres)}",
		0,
		PRIMARY_FONT_SIZE,
		False,
	)
	add_text_paragraph(
		text_frame,
		f"Director: {', '.join(movie_data.directors)}",
		0,
		PRIMARY_FONT_SIZE,
		False,
	)
	add_text_paragraph(
		text_frame,
		f"Run time: {movie_data.runtime_minutes} min",
		0,
		PRIMARY_FONT_SIZE,
		False,
	)
	add_text_paragraph(
		text_frame,
		f"Review Summary: {movie_data.rt_consensus}",
		0,
		PRIMARY_FONT_SIZE,
		False,
	)


#============================================
def place_poster(slide: object, movie_data: slide_maker.moviedata.MovieData) -> object:
	"""Replace the poster anchor with an aspect-preserving centered picture."""
	poster_anchor = shape_by_name(slide, TEMPLATE_POSTER_NAME)
	anchor_left = poster_anchor.left
	anchor_top = poster_anchor.top
	anchor_width = poster_anchor.width
	anchor_height = poster_anchor.height
	slide.shapes._spTree.remove(poster_anchor._element)

	with PIL.Image.open(movie_data.poster_path) as poster_image:
		poster_width, poster_height = poster_image.size
	require(poster_width > 0 and poster_height > 0, "Poster image has invalid dimensions")
	poster_ratio = poster_width / poster_height
	anchor_ratio = anchor_width / anchor_height
	if poster_ratio >= anchor_ratio:
		picture = slide.shapes.add_picture(
			movie_data.poster_path,
			anchor_left,
			anchor_top,
			width=anchor_width,
		)
		picture.top = anchor_top + (anchor_height - picture.height) // 2
	else:
		picture = slide.shapes.add_picture(
			movie_data.poster_path,
			anchor_left,
			anchor_top,
			height=anchor_height,
		)
		picture.left = anchor_left + (anchor_width - picture.width) // 2
	picture.name = TEMPLATE_POSTER_NAME
	picture._element.nvPicPr.cNvPr.set("descr", f"Poster for {movie_data.title}")
	return picture


#============================================
def append_movie_slide(presentation: object, movie_data: slide_maker.moviedata.MovieData) -> object:
	"""Append one visible movie slide using the presentation's template slide."""
	slide_maker.moviedata.validate_movie_data(movie_data)
	presentation.slide_width = PAGE_WIDTH
	presentation.slide_height = PAGE_HEIGHT
	require(len(presentation.slides) >= 1, "Presentation has no movie template slide")
	template_slide = presentation.slides[0]
	shape_by_name(template_slide, TEMPLATE_TITLE_NAME)
	shape_by_name(template_slide, TEMPLATE_OUTLINE_NAME)
	shape_by_name(template_slide, TEMPLATE_POSTER_NAME)

	new_slide = presentation.slides.add_slide(template_slide.slide_layout)
	clone_template_shapes(template_slide, new_slide)
	new_slide._element.attrib.pop("show", None)
	fill_title(new_slide, movie_data)
	fill_outline(new_slide, movie_data)
	place_poster(new_slide, movie_data)
	return new_slide


#============================================
def build_movie_presentation(
	movie_data: slide_maker.moviedata.MovieData,
	template_path: pathlib.Path,
	output_path: pathlib.Path,
) -> pathlib.Path:
	"""Load the extracted template, append one movie slide, and save the PPTX."""
	require(template_path.is_file(), f"Movie slide template is absent: {template_path}")
	presentation = pptx.Presentation(template_path)
	append_movie_slide(presentation, movie_data)
	output_path.parent.mkdir(parents=True, exist_ok=True)
	presentation.save(output_path)
	return output_path
