#!/usr/bin/env python3
"""Extract the authoritative hidden movie slide as the reusable template."""

# Standard Library
import pathlib
import zipfile
import argparse
import subprocess
import xml.etree.ElementTree

# PIP3 modules
import pptx
import pptx.enum.text


SOURCE_DECK = pathlib.Path("SLIDE_ARTIFACTS/class02b-pre-film_content.odp")
SOURCE_SLIDE_NUMBER = 2
OUTPUT_TEMPLATE = pathlib.Path("template/movie_slide_template.pptx")
EXPECTED_MASTER = "lecture99-template_2017"
EXPECTED_ODP_LAYOUT = "AL1T3"
EXPECTED_PAGE_WIDTH_CM = 28.0
EXPECTED_PAGE_HEIGHT_CM = 17.5
ANCHOR_NAMES = {
	"title": "movie_title",
	"outline": "movie_outline",
	"poster": "movie_poster",
}
REQUIRED_OUTLINE_LABELS = (
	"IMDB",
	"Critics: RT",
	"Genre:",
	"Director:",
	"Run time:",
	"Review Summary:",
)
NAMESPACES = {
	"draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
	"fo": "urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0",
	"presentation": "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0",
	"style": "urn:oasis:names:tc:opendocument:xmlns:style:1.0",
	"svg": "urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0",
}
DRAW_NAME = f"{{{NAMESPACES['draw']}}}name"
DRAW_STYLE_NAME = f"{{{NAMESPACES['draw']}}}style-name"
DRAW_MASTER_PAGE_NAME = f"{{{NAMESPACES['draw']}}}master-page-name"
FO_PAGE_WIDTH = f"{{{NAMESPACES['fo']}}}page-width"
FO_PAGE_HEIGHT = f"{{{NAMESPACES['fo']}}}page-height"
PRESENTATION_CLASS = f"{{{NAMESPACES['presentation']}}}class"
PRESENTATION_LAYOUT = f"{{{NAMESPACES['presentation']}}}presentation-page-layout-name"
PRESENTATION_VISIBILITY = f"{{{NAMESPACES['presentation']}}}visibility"
STYLE_DISPLAY_NAME = f"{{{NAMESPACES['style']}}}display-name"
STYLE_NAME = f"{{{NAMESPACES['style']}}}name"
STYLE_PAGE_LAYOUT_NAME = f"{{{NAMESPACES['style']}}}page-layout-name"
SVG_FONT_FAMILY = f"{{{NAMESPACES['svg']}}}font-family"


#============================================
def parse_args() -> argparse.Namespace:
	"""Parse the maintainer command line."""
	parser = argparse.ArgumentParser(
		description=(
			"Validate the ignored reference ODP and its companion PPTX, then replace "
			"the committed movie-slide template."
		),
	)
	args = parser.parse_args()
	return args


class TemplateExtractionError(RuntimeError):
	"""Report a missing source-deck structure required by the generator."""


#============================================
def get_repo_root() -> pathlib.Path:
	"""Return the repository root reported by Git."""
	result = subprocess.run(
		["git", "rev-parse", "--show-toplevel"],
		check=True,
		capture_output=True,
		text=True,
	)
	repo_root = pathlib.Path(result.stdout.strip())
	return repo_root


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise a template-specific error when a required condition is false."""
	if condition:
		return
	raise TemplateExtractionError(message)


#============================================
def read_package_xml(
	package_path: pathlib.Path,
	member_name: str,
) -> xml.etree.ElementTree.Element:
	"""Read one required XML member from an ODF package."""
	require(package_path.is_file(), f"Required source deck is absent: {package_path}")
	with zipfile.ZipFile(package_path) as archive:
		require(
			member_name in archive.namelist(),
			f"Required {member_name} is absent from source deck: {package_path}",
		)
		xml_bytes = archive.read(member_name)
	# The ignored reference deck is maintainer-owned local input.
	root = xml.etree.ElementTree.fromstring(xml_bytes)  # nosec B314
	return root


#============================================
def element_text(element: xml.etree.ElementTree.Element) -> str:
	"""Return normalized descendant text for one XML element."""
	text = " ".join("".join(element.itertext()).split())
	return text


#============================================
def require_one(
	elements: list[xml.etree.ElementTree.Element],
	description: str,
) -> xml.etree.ElementTree.Element:
	"""Return exactly one required XML element."""
	require(elements, f"Required {description} is absent")
	require(len(elements) == 1, f"Required {description} is ambiguous: found {len(elements)}")
	element = elements[0]
	return element


#============================================
def validate_odp_master_and_page(
	styles_root: xml.etree.ElementTree.Element,
) -> None:
	"""Validate the authoritative master, movie layout, page geometry, and font."""
	masters = [
		element
		for element in styles_root.findall(".//style:master-page", NAMESPACES)
		if element.get(STYLE_DISPLAY_NAME) == EXPECTED_MASTER
	]
	master = require_one(masters, f"master {EXPECTED_MASTER!r}")
	page_layout_name = master.get(STYLE_PAGE_LAYOUT_NAME)
	require(page_layout_name is not None, f"Master {EXPECTED_MASTER!r} has no page layout")

	page_layouts = [
		element
		for element in styles_root.findall(".//style:page-layout", NAMESPACES)
		if element.get(STYLE_NAME) == page_layout_name
	]
	page_layout = require_one(page_layouts, f"page layout {page_layout_name!r}")
	page_properties = require_one(
		page_layout.findall("style:page-layout-properties", NAMESPACES),
		f"page properties for {page_layout_name!r}",
	)
	require(
		page_properties.get(FO_PAGE_WIDTH) == "28cm"
		and page_properties.get(FO_PAGE_HEIGHT) == "17.5cm",
		"Required landscape page is not 28 cm by 17.5 cm",
	)

	presentation_layouts = [
		element
		for element in styles_root.findall(".//style:presentation-page-layout", NAMESPACES)
		if element.get(STYLE_NAME) == EXPECTED_ODP_LAYOUT
	]
	presentation_layout = require_one(
		presentation_layouts,
		f"movie presentation layout {EXPECTED_ODP_LAYOUT!r}",
	)
	layout_roles = [
		element.get(f"{{{NAMESPACES['presentation']}}}object")
		for element in presentation_layout.findall("presentation:placeholder", NAMESPACES)
	]
	require("title" in layout_roles, "Required title placeholder is absent from movie layout")
	require("outline" in layout_roles, "Required outline placeholder is absent from movie layout")

	master_roles = [
		element.get(PRESENTATION_CLASS)
		for element in master.findall("draw:frame", NAMESPACES)
	]
	require("title" in master_roles, "Required title placeholder is absent from movie master")
	require("outline" in master_roles, "Required outline placeholder is absent from movie master")

	font_families = {
		element.get(SVG_FONT_FAMILY)
		for element in styles_root.findall(".//style:font-face", NAMESPACES)
	}
	require("OpenDyslexic" in font_families, "Required OpenDyslexic font anchor is absent")


#============================================
def page_is_hidden(
	page: xml.etree.ElementTree.Element,
	content_root: xml.etree.ElementTree.Element,
) -> bool:
	"""Return whether the page's drawing style marks it hidden."""
	style_name = page.get(DRAW_STYLE_NAME)
	page_styles = [
		element
		for element in content_root.findall(".//style:style", NAMESPACES)
		if element.get(STYLE_NAME) == style_name
	]
	page_style = require_one(page_styles, f"drawing-page style {style_name!r}")
	properties = require_one(
		page_style.findall("style:drawing-page-properties", NAMESPACES),
		f"drawing-page properties for {style_name!r}",
	)
	is_hidden = properties.get(PRESENTATION_VISIBILITY) == "hidden"
	return is_hidden


#============================================
def select_odp_movie_page(
	content_root: xml.etree.ElementTree.Element,
) -> xml.etree.ElementTree.Element:
	"""Select and validate the hidden movie-template page in the source ODP."""
	candidates = []
	for page in content_root.findall(".//draw:page", NAMESPACES):
		page_text = element_text(page)
		if "TITLE (YEAR)" in page_text and "POSTER" in page_text:
			candidates.append(page)
	page = require_one(candidates, "hidden movie-template page")

	require(
		page.get(DRAW_NAME) == f"page{SOURCE_SLIDE_NUMBER}",
		f"Movie-template page is not expected slide {SOURCE_SLIDE_NUMBER}",
	)
	require(page_is_hidden(page, content_root), "Required movie-template page is visible")
	require(
		page.get(PRESENTATION_LAYOUT) == EXPECTED_ODP_LAYOUT,
		f"Movie-template page does not use layout {EXPECTED_ODP_LAYOUT!r}",
	)
	master_name = page.get(DRAW_MASTER_PAGE_NAME, "").replace("_5f_", "_")
	require(
		master_name == EXPECTED_MASTER,
		f"Movie-template page does not use master {EXPECTED_MASTER!r}",
	)

	frames = page.findall("draw:frame", NAMESPACES)
	title_frames = [frame for frame in frames if frame.get(PRESENTATION_CLASS) == "title"]
	outline_frames = [frame for frame in frames if frame.get(PRESENTATION_CLASS) == "outline"]
	title_frame = require_one(title_frames, "movie title anchor")
	outline_frame = require_one(outline_frames, "movie outline anchor")
	require("TITLE (YEAR)" in element_text(title_frame), "Movie title anchor has no title marker")
	outline_text = element_text(outline_frame)
	for label in REQUIRED_OUTLINE_LABELS:
		require(label in outline_text, f"Movie outline anchor is missing label {label!r}")

	poster_candidates = [
		element
		for element in page.findall("draw:custom-shape", NAMESPACES)
		if element_text(element) == "POSTER"
	]
	require_one(poster_candidates, "movie poster anchor")
	page_anchors = [
		element
		for element in page.findall(".//draw:page-thumbnail", NAMESPACES)
		if element.get(PRESENTATION_CLASS) == "page"
	]
	require_one(page_anchors, "notes page anchor")
	return page


#============================================
def shape_text(shape: object) -> str:
	"""Return normalized text from a presentation shape."""
	text = getattr(shape, "text", "")
	normalized = " ".join(text.split())
	return normalized


#============================================
def require_shape(candidates: list[object], description: str) -> object:
	"""Return exactly one required presentation shape."""
	require(candidates, f"Required {description} is absent from companion PPTX")
	require(
		len(candidates) == 1,
		f"Required {description} is ambiguous in companion PPTX: found {len(candidates)}",
	)
	shape = candidates[0]
	return shape


#============================================
def first_run_font(shape: object) -> str | None:
	"""Return the first explicit run font from a text shape."""
	for paragraph in shape.text_frame.paragraphs:
		for run in paragraph.runs:
			if run.font.name is not None:
				font_name = run.font.name
				return font_name
	return None


#============================================
def find_pptx_anchors(slide: object) -> dict[str, object]:
	"""Find and validate semantic anchors on the companion movie slide."""
	shapes = list(slide.shapes)
	title_candidates = [
		shape for shape in shapes
		if getattr(shape, "is_placeholder", False) and "TITLE (YEAR)" in shape_text(shape)
	]
	outline_candidates = [
		shape for shape in shapes
		if getattr(shape, "is_placeholder", False)
		and all(label in shape_text(shape) for label in REQUIRED_OUTLINE_LABELS)
	]
	poster_candidates = [
		shape for shape in shapes
		if shape_text(shape) == "POSTER" and not getattr(shape, "is_placeholder", False)
	]
	anchors = {
		"title": require_shape(title_candidates, "movie title anchor"),
		"outline": require_shape(outline_candidates, "movie outline anchor"),
		"poster": require_shape(poster_candidates, "movie poster anchor"),
	}

	for role, shape in anchors.items():
		require(
			shape.width > 0 and shape.height > 0,
			f"Required {role} anchor has invalid geometry",
		)
		require(
			first_run_font(shape) == "OpenDyslexic",
			f"Required {role} font anchor is not OpenDyslexic",
		)
	require(
		anchors["title"].text_frame.auto_size == pptx.enum.text.MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
		"Required title autofit anchor is absent",
	)
	require(
		anchors["outline"].text_frame.auto_size == pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
		"Required outline autofit anchor is absent",
	)
	return anchors


#============================================
def select_pptx_movie_slide(presentation: object) -> tuple[int, object, dict[str, object]]:
	"""Select the companion PPTX slide matching the authoritative ODP page."""
	candidates = []
	for index, slide in enumerate(presentation.slides):
		slide_text = " ".join(shape_text(shape) for shape in slide.shapes)
		if "TITLE (YEAR)" in slide_text and "POSTER" in slide_text:
			candidates.append((index, slide))
	require(candidates, "Required movie-template slide is absent from companion PPTX")
	require(
		len(candidates) == 1,
		f"Movie-template slide is ambiguous in companion PPTX: found {len(candidates)}",
	)
	slide_index, slide = candidates[0]
	require(
		slide_index + 1 == SOURCE_SLIDE_NUMBER,
		f"Companion movie-template slide is not slide {SOURCE_SLIDE_NUMBER}",
	)
	require(slide._element.get("show") == "0", "Companion movie-template slide is visible")
	require(
		slide.slide_layout.name == EXPECTED_MASTER,
		f"Required companion layout {EXPECTED_MASTER!r} is absent",
	)
	anchors = find_pptx_anchors(slide)
	return slide_index, slide, anchors


#============================================
def validate_pptx_structure(presentation: object) -> None:
	"""Validate required master, layout, placeholder, and page anchors."""
	require(len(presentation.slide_masters) == 1, "Companion PPTX does not have one master")
	master = presentation.slide_masters[0]
	layouts = [layout for layout in master.slide_layouts if layout.name == EXPECTED_MASTER]
	layout = require_shape(layouts, f"layout {EXPECTED_MASTER!r}")
	require(len(layout.placeholders) >= 2, "Required layout placeholders are absent")

	width_delta = abs(presentation.slide_width - int(EXPECTED_PAGE_WIDTH_CM * 360000))
	height_delta = abs(presentation.slide_height - int(EXPECTED_PAGE_HEIGHT_CM * 360000))
	require(width_delta <= 2000, "Companion PPTX page width does not match 28 cm")
	require(height_delta <= 2000, "Companion PPTX page height does not match 17.5 cm")


#============================================
def retain_selected_slide(presentation: object, selected_index: int) -> None:
	"""Remove every slide except the selected hidden movie-template slide."""
	slide_ids = presentation.slides._sldIdLst
	for index in reversed(range(len(slide_ids))):
		if index == selected_index:
			continue
		relationship_id = slide_ids[index].rId
		presentation.part.drop_rel(relationship_id)
		del slide_ids[index]


#============================================
def record_template_identity(presentation: object) -> None:
	"""Record the authoritative source selection in package metadata."""
	presentation.core_properties.title = "Movie slide template"
	presentation.core_properties.subject = (
		f"Extracted from {SOURCE_DECK.as_posix()} slide {SOURCE_SLIDE_NUMBER}"
	)
	presentation.core_properties.keywords = (
		f"master={EXPECTED_MASTER}; layout={EXPECTED_ODP_LAYOUT}; "
		"anchors=movie_title,movie_outline,movie_poster"
	)


#============================================
def write_template(
	source_pptx: pathlib.Path,
	output_path: pathlib.Path,
) -> dict[str, object]:
	"""Extract the validated hidden movie slide into a one-slide PPTX."""
	require(source_pptx.is_file(), f"Required companion PPTX is absent: {source_pptx}")
	presentation = pptx.Presentation(source_pptx)
	validate_pptx_structure(presentation)
	slide_index, slide, anchors = select_pptx_movie_slide(presentation)
	for role, shape in anchors.items():
		shape.name = ANCHOR_NAMES[role]
	retain_selected_slide(presentation, slide_index)
	record_template_identity(presentation)
	output_path.parent.mkdir(parents=True, exist_ok=True)
	presentation.save(output_path)
	return anchors


#============================================
def validate_written_template(output_path: pathlib.Path) -> dict[str, object]:
	"""Reopen the generated template and verify its durable anchor map."""
	require(output_path.is_file(), f"Generated template is absent: {output_path}")
	presentation = pptx.Presentation(output_path)
	validate_pptx_structure(presentation)
	require(len(presentation.slides) == 1, "Generated template does not contain one slide")
	slide = presentation.slides[0]
	require(slide._element.get("show") == "0", "Generated template slide is visible")
	anchors = find_pptx_anchors(slide)
	for role, shape in anchors.items():
		require(
			shape.name == ANCHOR_NAMES[role],
			f"Generated template {role} anchor name was not recorded",
		)
	expected_subject = f"Extracted from {SOURCE_DECK.as_posix()} slide {SOURCE_SLIDE_NUMBER}"
	require(
		presentation.core_properties.subject == expected_subject,
		"Generated template source deck and slide were not recorded",
	)
	return anchors


#============================================
def describe_anchor(role: str, shape: object) -> str:
	"""Format one discovered anchor for maintainer inspection."""
	left_cm = shape.left / 360000
	top_cm = shape.top / 360000
	width_cm = shape.width / 360000
	height_cm = shape.height / 360000
	auto_size = shape.text_frame.auto_size.name
	description = (
		f"{role}: name={shape.name}, x={left_cm:.2f}cm, y={top_cm:.2f}cm, "
		f"w={width_cm:.2f}cm, h={height_cm:.2f}cm, "
		f"font={first_run_font(shape)}, autofit={auto_size}"
	)
	return description


#============================================
def main() -> None:
	"""Validate the source and write the extracted movie-slide template."""
	parse_args()
	repo_root = get_repo_root()
	source_deck = repo_root / SOURCE_DECK
	source_pptx = source_deck.with_suffix(".pptx")
	output_path = repo_root / OUTPUT_TEMPLATE

	content_root = read_package_xml(source_deck, "content.xml")
	styles_root = read_package_xml(source_deck, "styles.xml")
	validate_odp_master_and_page(styles_root)
	page = select_odp_movie_page(content_root)
	write_template(source_pptx, output_path)
	anchors = validate_written_template(output_path)

	print(f"Selected source deck: {SOURCE_DECK.as_posix()}")
	print(f"Selected source slide: {SOURCE_SLIDE_NUMBER} ({page.get(DRAW_NAME)}, hidden)")
	print(f"Selected master: {EXPECTED_MASTER}")
	print(f"Selected movie layout: {EXPECTED_ODP_LAYOUT}")
	print(f"Page anchor: {EXPECTED_PAGE_WIDTH_CM:g}cm x {EXPECTED_PAGE_HEIGHT_CM:g}cm landscape")
	for role in ("title", "outline", "poster"):
		print(describe_anchor(role, anchors[role]))
	print(f"Wrote template: {OUTPUT_TEMPLATE.as_posix()}")


if __name__ == "__main__":
	main()
