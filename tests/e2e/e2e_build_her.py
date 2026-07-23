"""Build and inspect a Her movie slide using the real extracted template."""

# Standard Library
import sys
import pathlib
import tempfile

# PIP3 modules
import pptx
import PIL.Image #pillow
import pptx.oxml.ns
import pptx.enum.text
import pptx.enum.shapes


TESTS_DIR = pathlib.Path(__file__).resolve().parents[1]
sys.path.insert(0, str(TESTS_DIR))

# local repo modules
import file_utils

REPO_ROOT = pathlib.Path(file_utils.get_repo_root())

# local repo modules
import slide_maker.moviedata
import slide_maker.slide_builder


EXPECTED_LABELS = (
	"IMDB rating",
	"Critics: RT",
	"Audience:",
	"Genre:",
	"Director:",
	"Run time:",
	"Review Summary:",
)


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise a clear E2E failure when a required semantic property is absent."""
	if condition:
		return
	raise RuntimeError(message)


#============================================
def shape_by_name(slide: object, name: str) -> object:
	"""Return one named shape without relying on package element order."""
	candidates = [shape for shape in slide.shapes if shape.name == name]
	require(candidates, f"Built slide is missing semantic role: {name}")
	require(len(candidates) == 1, f"Built slide has ambiguous semantic role: {name}")
	shape = candidates[0]
	return shape


#============================================
def find_built_slide(presentation: object) -> object:
	"""Find the visible Her slide by semantic title content."""
	candidates = []
	for slide in presentation.slides:
		if slide._element.get("show") == "0":
			continue
		title = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_TITLE_NAME)
		if "Her (2013)" in title.text:
			candidates.append(slide)
	require(candidates, "Built Her slide is absent")
	require(len(candidates) == 1, "Built Her slide is ambiguous")
	slide = candidates[0]
	return slide


#============================================
def verify_text_roles(slide: object) -> None:
	"""Verify literal labels, hierarchy, font, and autofit by semantic role."""
	title = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_TITLE_NAME)
	outline = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME)
	require(title.text.strip() == "Her (2013)", "Built title does not match Her identity")
	require(
		title.text_frame.auto_size == pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
		"Built title is missing text-to-fit autofit",
	)
	require(
		outline.text_frame.auto_size == pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
		"Built outline is missing text-to-fit autofit",
	)
	for label in EXPECTED_LABELS:
		require(label in outline.text, f"Built outline is missing literal label: {label}")
	levels = [paragraph.level for paragraph in outline.text_frame.paragraphs]
	require(levels == [0, 1, 1, 1, 0, 0, 0, 0], "Built outline hierarchy changed")
	for shape in (title, outline):
		for paragraph in shape.text_frame.paragraphs:
			for run in paragraph.runs:
				require(run.font.name == "OpenDyslexic", "Built text font changed")
	imdb_paragraph = outline.text_frame.paragraphs[1]
	score_runs = [run for run in imdb_paragraph.runs if run.text == "8.0"]
	require(len(score_runs) == 1, "Built IMDb score run is absent or ambiguous")
	run_properties = score_runs[0]._r.get_or_add_rPr()
	highlight = run_properties.find(pptx.oxml.ns.qn("a:highlight"))
	require(highlight is not None, "Built IMDb score has no color highlight")
	color = highlight.find(pptx.oxml.ns.qn("a:srgbClr"))
	require(color is not None, "Built IMDb score highlight has no explicit RGB color")
	require(
		color.get("val") == slide_maker.slide_builder.IMDB_SCORE_HIGHLIGHT,
		"Built IMDb score highlight color changed",
	)
	require("792k votes" in imdb_paragraph.text, "Built IMDb vote count is not compact")
	ratings_text = outline.text_frame.paragraphs[2].text
	require(
		slide_maker.emoji_marks.rt_critic_mark_for_score(95) in ratings_text,
		"Built critics score has no high-score mark",
	)
	require(
		slide_maker.emoji_marks.rt_audience_mark_for_score(82)
		in outline.text_frame.paragraphs[3].text,
		"Built audience score has no Popcornmeter mark",
	)


#============================================
def verify_compact_count_examples() -> None:
	"""Verify the requested compact count boundaries and precision."""
	examples = {
		435_444: "435k",
		1_200_000: "1.2M",
		2_300: "2.3k",
		1_000: "1.0k",
		999_499: "999k",
		999_500: "1.0M",
	}
	for count, expected in examples.items():
		actual = slide_maker.slide_builder.format_compact_count(count)
		require(actual == expected, f"Compact count {count} rendered as {actual}, expected {expected}")


#============================================
def verify_rating_mark_boundaries() -> None:
	"""Verify the requested Rotten Tomatoes critic and audience tiers."""
	critic_examples = {
		59: slide_maker.emoji_marks.ROTTEN_MARK,
		60: slide_maker.emoji_marks.TOMATO_MARK,
		80: slide_maker.emoji_marks.TOMATO_MARK,
		81: (
			slide_maker.emoji_marks.TOMATO_MARK
			+ slide_maker.emoji_marks.TROPHY_MARK
		),
	}
	for score, expected in critic_examples.items():
		actual = slide_maker.emoji_marks.rt_critic_mark_for_score(score)
		require(actual == expected, f"Critics score {score} used the wrong display mark")
	audience_examples = {
		59: (
			slide_maker.emoji_marks.POPCORN_MARK
			+ slide_maker.emoji_marks.THUMBS_DOWN_MARK
		),
		60: slide_maker.emoji_marks.POPCORN_MARK,
	}
	for score, expected in audience_examples.items():
		actual = slide_maker.emoji_marks.rt_audience_mark_for_score(score)
		require(actual == expected, f"Audience score {score} used the wrong display mark")


#============================================
def verify_poster_role(slide: object, template_slide: object) -> None:
	"""Verify poster containment, centering, and preserved source aspect."""
	reference = shape_by_name(template_slide, slide_maker.slide_builder.TEMPLATE_POSTER_NAME)
	poster = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_POSTER_NAME)
	require(
		poster.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE,
		"Built poster role is not a picture",
	)
	require(poster.left >= reference.left, "Built poster extends left of its reference region")
	require(poster.top >= reference.top, "Built poster extends above its reference region")
	require(
		poster.left + poster.width <= reference.left + reference.width,
		"Built poster extends right of its reference region",
	)
	require(
		poster.top + poster.height <= reference.top + reference.height,
		"Built poster extends below its reference region",
	)
	reference_center_x = reference.left + reference.width / 2
	reference_center_y = reference.top + reference.height / 2
	poster_center_x = poster.left + poster.width / 2
	poster_center_y = poster.top + poster.height / 2
	require(abs(reference_center_x - poster_center_x) <= 1, "Built poster is not centered")
	require(abs(reference_center_y - poster_center_y) <= 1, "Built poster is not centered")
	poster_ratio = poster.width / poster.height
	require(abs(poster_ratio - (2 / 3)) < 0.001, "Built poster aspect ratio changed")


#============================================
def main() -> None:
	"""Generate a runtime poster, build a real PPTX, and inspect its semantics."""
	template_path = REPO_ROOT / "template" / "movie_slide_template.pptx"
	with tempfile.TemporaryDirectory() as temporary_directory:
		work_dir = pathlib.Path(temporary_directory)
		poster_path = work_dir / "her_poster.png"
		output_path = work_dir / "her_slide.pptx"
		poster_image = PIL.Image.new("RGB", (600, 900), color=(35, 72, 115))
		poster_image.save(poster_path)
		movie_data = slide_maker.moviedata.MovieData(
			title="Her",
			year=2013,
			plot="A lonely writer develops an unexpected relationship with an operating system.",
			genres=["Drama", "Romance", "Science Fiction"],
			runtime_minutes=126,
			directors=["Spike Jonze"],
			tmdb_id=152601,
			imdb_id="tt1798709",
			imdb_rating=8.0,
			imdb_votes=792315,
			rt_tomatometer=95,
			rt_audience_score=82,
			rt_state="fresh",
			rt_consensus="Sweet, soulful, and smart, Her uses its scenario to impart wisdom.",
			metascore=91,
			metascore_band="high",
			poster_path=str(poster_path),
		)
		slide_maker.slide_builder.build_movie_presentation(
			movie_data,
			template_path,
			output_path,
		)
		require(output_path.is_file(), "Slide builder did not write the presentation")
		presentation = pptx.Presentation(output_path)
		require(
			presentation.slide_width == slide_maker.slide_builder.PAGE_WIDTH
			and presentation.slide_height == slide_maker.slide_builder.PAGE_HEIGHT,
			"Built presentation does not use the exact lecture page size",
		)
		template_slide = presentation.slides[0]
		require(template_slide._element.get("show") == "0", "Template slide became visible")
		built_slide = find_built_slide(presentation)
		verify_text_roles(built_slide)
		verify_poster_role(built_slide, template_slide)
		verify_compact_count_examples()
		verify_rating_mark_boundaries()
		print(f"Slide-builder E2E passed: {output_path}")


if __name__ == "__main__":
	main()
