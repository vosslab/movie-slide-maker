"""Accept normal and long-text movie slides through real ODP rendering."""

# Standard Library
import sys
import pathlib
import zipfile
import tempfile
import subprocess
import xml.etree.ElementTree

# PIP3 modules
import pptx
import PIL.Image #pillow
import PIL.ImageStat #pillow
import PIL.ImageChops #pillow
import pptx.enum.text


TESTS_DIR = pathlib.Path(__file__).resolve().parents[1]
sys.path.insert(0, str(TESTS_DIR))

# local repo modules
import file_utils

REPO_ROOT = pathlib.Path(file_utils.get_repo_root())

# local repo modules
import slide_maker.moviedata
import slide_maker.emoji_marks
import slide_maker.libreoffice_runner
import slide_maker.slide_builder
import slide_maker.slide_convert


DRAW_NS = slide_maker.slide_convert.DRAW_NS
FO_NS = slide_maker.slide_convert.FO_NS
OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
PRESENTATION_NS = slide_maker.slide_convert.PRESENTATION_NS
STYLE_NS = slide_maker.slide_convert.STYLE_NS
SVG_NS = slide_maker.slide_convert.SVG_NS
TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
EXPECTED_PAGE_WIDTH_CM = 28.0
EXPECTED_PAGE_HEIGHT_CM = 17.5
GEOMETRY_TOLERANCE_CM = slide_maker.slide_convert.GEOMETRY_TOLERANCE_CM
EVIDENCE_DIR = REPO_ROOT / "output_smoke" / "visual_accept"
EXPECTED_LABELS = (
	"IMDB rating",
	"Critics: RT",
	"Audience:",
	"Genre:",
	"Director:",
	"Run time:",
	"Review Summary:",
)
GEOMETRY_ATTRIBUTES = (
	("x", "left"),
	("y", "top"),
	("width", "width"),
	("height", "height"),
)


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise a clear visual-acceptance failure when a condition is false."""
	if condition:
		return
	raise RuntimeError(message)


#============================================
def normalized_text(element: xml.etree.ElementTree.Element) -> str:
	"""Return normalized descendant text for stable semantic comparisons."""
	text = " ".join("".join(element.itertext()).split())
	return text


#============================================
def shape_by_name(slide: object, name: str) -> object:
	"""Return one named source shape without depending on element order."""
	candidates = [shape for shape in slide.shapes if shape.name == name]
	require(candidates, f"Built slide is missing semantic role: {name}")
	require(len(candidates) == 1, f"Built slide has ambiguous semantic role: {name}")
	shape = candidates[0]
	return shape


#============================================
def cm_geometry(shape: object) -> dict[str, float]:
	"""Return source shape geometry converted from EMUs to centimeters."""
	geometry = {
		"x": float(shape.left) / 360000.0,
		"y": float(shape.top) / 360000.0,
		"width": float(shape.width) / 360000.0,
		"height": float(shape.height) / 360000.0,
	}
	return geometry


#============================================
def movie_cases(work_dir: pathlib.Path) -> tuple[slide_maker.moviedata.MovieData, ...]:
	"""Create normal and realistic long-text movie records with runtime posters."""
	normal_poster_path = work_dir / "normal_poster.png"
	long_poster_path = work_dir / "long_text_poster.png"
	PIL.Image.new("RGB", (600, 900), color=(35, 72, 115)).save(normal_poster_path)
	PIL.Image.new("RGB", (800, 1200), color=(116, 55, 38)).save(long_poster_path)
	normal = slide_maker.moviedata.MovieData(
		title="Her",
		year=2013,
		plot="A lonely writer develops an unexpected relationship with an operating system.",
		genres=["Drama", "Romance", "Science Fiction"],
		runtime_minutes=126,
		directors=["Spike Jonze"],
		tmdb_id=152601,
		imdb_id="tt1798709",
		imdb_rating=8.0,
		imdb_votes=792315,
		rt_tomatometer=95,
		rt_audience_score=82,
		rt_state="fresh",
		rt_consensus="Sweet, soulful, and smart, Her uses its scenario to impart wisdom.",
		metascore=91,
		metascore_band="high",
		poster_path=str(normal_poster_path),
	)
	long_text = slide_maker.moviedata.MovieData(
		title="Her",
		year=2013,
		plot=(
			"In a near-future Los Angeles, a lonely writer recovering from the end of his marriage "
			"develops an unexpected relationship with a perceptive operating system, then confronts "
			"how intimacy, memory, creativity, and personal growth can remain meaningful when the "
			"person he loves experiences the world in a radically different way."
		),
		genres=["Drama", "Romance", "Science Fiction"],
		runtime_minutes=126,
		directors=["Spike Jonze"],
		tmdb_id=152601,
		imdb_id="tt1798709",
		imdb_rating=8.0,
		imdb_votes=792315,
		rt_tomatometer=95,
		rt_audience_score=82,
		rt_state="fresh",
		rt_consensus=(
			"Sweet, soulful, visually thoughtful, and anchored by a warm central performance, Her "
			"uses its speculative relationship to explore loneliness, emotional change, and the "
			"difficult wisdom required to let another person grow beyond our expectations."
		),
		metascore=91,
		metascore_band="high",
		poster_path=str(long_poster_path),
	)
	return normal, long_text


#============================================
def source_movie_slide(presentation: object, expected_title: str) -> object:
	"""Return the one visible source slide with the expected movie title."""
	candidates = []
	for slide in presentation.slides:
		if slide._element.get("show") == "0":
			continue
		title = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_TITLE_NAME)
		if title.text.strip() == expected_title:
			candidates.append(slide)
	require(candidates, f"Built presentation is missing movie slide: {expected_title}")
	require(len(candidates) == 1, f"Built movie slide is ambiguous: {expected_title}")
	slide = candidates[0]
	return slide


#============================================
def expected_outline_fragments(movie_data: slide_maker.moviedata.MovieData) -> tuple[str, ...]:
	"""Return stable literal movie fragments required in the outline."""
	fragments = [
		movie_data.plot,
		f"{movie_data.imdb_rating:.1f}",
		slide_maker.slide_builder.format_compact_count(movie_data.imdb_votes),
		f"{movie_data.rt_tomatometer}%",
		str(movie_data.metascore),
		", ".join(movie_data.genres),
		", ".join(movie_data.directors),
		f"{movie_data.runtime_minutes} min",
		movie_data.rt_consensus,
		slide_maker.emoji_marks.rt_critic_mark_for_score(movie_data.rt_tomatometer),
		slide_maker.emoji_marks.GREEN_SQUARE_MARK,
	]
	if movie_data.rt_audience_score is None:
		fragments.append("Audience: N/A")
	else:
		fragments.extend(
			(
				f"{movie_data.rt_audience_score}%",
				slide_maker.emoji_marks.rt_audience_mark_for_score(
					movie_data.rt_audience_score
				),
			)
		)
	return tuple(fragments)


#============================================
def validate_source_text(slide: object, movie_data: slide_maker.moviedata.MovieData) -> None:
	"""Validate source roles, labels, bullet meaning, font, and autofit."""
	title = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_TITLE_NAME)
	outline = shape_by_name(slide, slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME)
	expected_title = f"{movie_data.title} ({movie_data.year})"
	require(title.text.strip() == expected_title, "Built title changed movie identity")
	require(
		title.text_frame.auto_size == pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
		"Built title is missing text-to-fit autofit",
	)
	require(
		outline.text_frame.auto_size == pptx.enum.text.MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
		"Built outline is missing text-to-fit autofit",
	)
	for label in EXPECTED_LABELS:
		require(label in outline.text, f"Built outline is missing literal label: {label}")
	for fragment in expected_outline_fragments(movie_data):
		require(fragment in outline.text, f"Built outline is missing movie text: {fragment}")
	paragraphs = outline.text_frame.paragraphs
	levels = [paragraph.level for paragraph in paragraphs]
	require(levels == [0, 1, 1, 1, 0, 0, 0, 0], "Built outline bullet hierarchy changed")
	require("IMDB rating" in paragraphs[1].text, "IMDb rating lost its secondary bullet")
	require("Critics: RT" in paragraphs[2].text, "Critics ratings lost their secondary bullet")
	require("Audience:" in paragraphs[3].text, "Popcornmeter lost its secondary bullet")
	for shape in (title, outline):
		for paragraph in shape.text_frame.paragraphs:
			for run in paragraph.runs:
				require(
					run.font.name == slide_maker.slide_builder.FONT_NAME,
					f"Built role {shape.name} is not OpenDyslexic",
				)


#============================================
def validate_source_poster(
	slide: object,
	template_slide: object,
	movie_data: slide_maker.moviedata.MovieData,
) -> None:
	"""Validate source poster containment, centering, and aspect preservation."""
	role = slide_maker.slide_builder.TEMPLATE_POSTER_NAME
	poster = shape_by_name(slide, role)
	anchor = shape_by_name(template_slide, role)
	require(poster.left >= anchor.left, "Built poster extends left of its intended region")
	require(poster.top >= anchor.top, "Built poster extends above its intended region")
	require(
		poster.left + poster.width <= anchor.left + anchor.width,
		"Built poster extends right of its intended region",
	)
	require(
		poster.top + poster.height <= anchor.top + anchor.height,
		"Built poster extends below its intended region",
	)
	anchor_center_x = anchor.left + anchor.width / 2
	anchor_center_y = anchor.top + anchor.height / 2
	poster_center_x = poster.left + poster.width / 2
	poster_center_y = poster.top + poster.height / 2
	require(abs(anchor_center_x - poster_center_x) <= 1, "Built poster is not centered horizontally")
	require(abs(anchor_center_y - poster_center_y) <= 1, "Built poster is not centered vertically")
	with PIL.Image.open(movie_data.poster_path) as poster_image:
		image_ratio = poster_image.width / poster_image.height
	frame_ratio = poster.width / poster.height
	require(abs(frame_ratio - image_ratio) < 0.001, "Built poster aspect ratio changed")


#============================================
def build_and_validate_source(
	movie_data: slide_maker.moviedata.MovieData,
	template_path: pathlib.Path,
	scratch_path: pathlib.Path,
) -> tuple[dict[str, dict[str, float]], dict[str, float]]:
	"""Build a real PPTX and return accepted role and poster-anchor geometry."""
	slide_maker.slide_builder.build_movie_presentation(movie_data, template_path, scratch_path)
	require(scratch_path.is_file(), f"Builder did not create scratch presentation: {scratch_path}")
	presentation = pptx.Presentation(scratch_path)
	template_slide = presentation.slides[0]
	require(template_slide._element.get("show") == "0", "Template slide became visible")
	expected_title = f"{movie_data.title} ({movie_data.year})"
	slide = source_movie_slide(presentation, expected_title)
	validate_source_text(slide, movie_data)
	validate_source_poster(slide, template_slide, movie_data)
	role_geometry = {}
	for role in slide_maker.slide_convert.REQUIRED_ROLES:
		role_geometry[role] = cm_geometry(shape_by_name(slide, role))
	poster_anchor = cm_geometry(
		shape_by_name(template_slide, slide_maker.slide_builder.TEMPLATE_POSTER_NAME)
	)
	return role_geometry, poster_anchor


#============================================
def write_empty_text_control(
	scratch_path: pathlib.Path,
	control_path: pathlib.Path,
	movie_data: slide_maker.moviedata.MovieData,
	role_name: str,
	role_label: str,
) -> None:
	"""Write an otherwise identical deck with one semantic text role cleared."""
	presentation = pptx.Presentation(scratch_path)
	expected_title = f"{movie_data.title} ({movie_data.year})"
	slide = source_movie_slide(presentation, expected_title)
	text_shape = shape_by_name(slide, role_name)
	text_shape.text_frame.clear()
	presentation.save(control_path)
	require(control_path.is_file(), f"Empty-{role_label} control was not written: {control_path}")


#============================================
def parse_odp(
	odp_path: pathlib.Path,
) -> tuple[
	zipfile.ZipFile,
	xml.etree.ElementTree.Element,
	xml.etree.ElementTree.Element,
]:
	"""Open an ODP and parse the trusted LibreOffice XML members."""
	require(zipfile.is_zipfile(odp_path), f"Converted output is not an ODP package: {odp_path}")
	archive = zipfile.ZipFile(odp_path)
	content_root = slide_maker.slide_convert.parse_xml(archive, "content.xml")
	styles_root = slide_maker.slide_convert.parse_xml(archive, "styles.xml")
	return archive, content_root, styles_root


#============================================
def page_dimensions_cm(
	page: xml.etree.ElementTree.Element,
	styles_root: xml.etree.ElementTree.Element,
) -> tuple[float, float]:
	"""Resolve converted page dimensions through its master and page layout."""
	master_name = page.get(f"{{{DRAW_NS}}}master-page-name")
	require(bool(master_name), "Converted movie slide has no master-page reference")
	masters = [
		master for master in styles_root.iter(f"{{{STYLE_NS}}}master-page")
		if master.get(f"{{{STYLE_NS}}}name") == master_name
	]
	require(len(masters) == 1, "Converted movie slide master-page reference is invalid")
	layout_name = masters[0].get(f"{{{STYLE_NS}}}page-layout-name")
	require(bool(layout_name), "Converted movie slide master has no page layout")
	layouts = [
		layout for layout in styles_root.iter(f"{{{STYLE_NS}}}page-layout")
		if layout.get(f"{{{STYLE_NS}}}name") == layout_name
	]
	require(len(layouts) == 1, "Converted movie slide page-layout reference is invalid")
	properties = layouts[0].find(f"{{{STYLE_NS}}}page-layout-properties")
	require(properties is not None, "Converted movie slide page layout has no properties")
	width_value = properties.get(f"{{{FO_NS}}}page-width")
	height_value = properties.get(f"{{{FO_NS}}}page-height")
	require(bool(width_value and height_value), "Converted movie slide has no page dimensions")
	width_cm = slide_maker.slide_convert.length_cm(width_value)
	height_cm = slide_maker.slide_convert.length_cm(height_value)
	return width_cm, height_cm


#============================================
def validate_page(
	content_root: xml.etree.ElementTree.Element,
	styles_root: xml.etree.ElementTree.Element,
) -> xml.etree.ElementTree.Element:
	"""Validate one visible landscape page against the reference page geometry."""
	presentation = content_root.find(f".//{{{OFFICE_NS}}}presentation")
	require(presentation is not None, "Converted ODP has no presentation body")
	visibility_attribute = f"{{{PRESENTATION_NS}}}visibility"
	style_name_attribute = f"{{{DRAW_NS}}}style-name"
	style_name_map = {
		style.get(f"{{{STYLE_NS}}}name"): style
		for style in content_root.iter(f"{{{STYLE_NS}}}style")
	}
	visible_pages = []
	for page_candidate in presentation.findall(f"{{{DRAW_NS}}}page"):
		page_style = style_name_map.get(page_candidate.get(style_name_attribute))
		drawing_properties = None
		if page_style is not None:
			drawing_properties = page_style.find(f"{{{STYLE_NS}}}drawing-page-properties")
		visibility = None
		if drawing_properties is not None:
			visibility = drawing_properties.get(visibility_attribute)
		if visibility != "hidden":
			visible_pages.append(page_candidate)
	require(len(visible_pages) == 1, "Converted deck does not have exactly one visible page")
	page = slide_maker.slide_convert.movie_pages(content_root, 1)[0]
	require(page is visible_pages[0], "Visible page does not carry the movie semantic roles")
	width_cm, height_cm = page_dimensions_cm(page, styles_root)
	require(width_cm > height_cm, "Converted movie slide is not landscape")
	require(
		abs(width_cm - EXPECTED_PAGE_WIDTH_CM) <= GEOMETRY_TOLERANCE_CM,
		f"Converted page width changed: {width_cm:.3f} cm",
	)
	require(
		abs(height_cm - EXPECTED_PAGE_HEIGHT_CM) <= GEOMETRY_TOLERANCE_CM,
		f"Converted page height changed: {height_cm:.3f} cm",
	)
	expected_ratio = EXPECTED_PAGE_WIDTH_CM / EXPECTED_PAGE_HEIGHT_CM
	require(abs((width_cm / height_cm) - expected_ratio) < 0.01, "Converted page aspect changed")
	return page


#============================================
def style_map(
	content_root: xml.etree.ElementTree.Element,
	styles_root: xml.etree.ElementTree.Element,
) -> dict[str, xml.etree.ElementTree.Element]:
	"""Collect named ODF styles from content and styles XML."""
	styles = {}
	name_attribute = f"{{{STYLE_NS}}}name"
	for root in (content_root, styles_root):
		for style in root.iter(f"{{{STYLE_NS}}}style"):
			name = style.get(name_attribute)
			if name:
				styles[name] = style
	return styles


#============================================
def referenced_style_values(
	frame: xml.etree.ElementTree.Element,
	styles: dict[str, xml.etree.ElementTree.Element],
) -> set[str]:
	"""Resolve property values from styles referenced by a semantic frame."""
	style_names = set()
	for element in frame.iter():
		for attribute_name, value in element.attrib.items():
			if attribute_name.endswith("}style-name"):
				style_names.add(value)
	values = set()
	pending = list(style_names)
	visited = set()
	parent_attribute = f"{{{STYLE_NS}}}parent-style-name"
	while pending:
		style_name = pending.pop()
		if style_name in visited:
			continue
		visited.add(style_name)
		style = styles.get(style_name)
		if style is None:
			continue
		values.update(style.attrib.values())
		for descendant in style.iter():
			values.update(descendant.attrib.values())
		parent_name = style.get(parent_attribute)
		if parent_name:
			pending.append(parent_name)
	return values


#============================================
def validate_converted_text(
	page: xml.etree.ElementTree.Element,
	content_root: xml.etree.ElementTree.Element,
	styles_root: xml.etree.ElementTree.Element,
	movie_data: slide_maker.moviedata.MovieData,
) -> None:
	"""Validate converted labels, movie meaning, fonts, bullets, and autofit."""
	frames = slide_maker.slide_convert.named_frames(page)
	title_role = slide_maker.slide_builder.TEMPLATE_TITLE_NAME
	outline_role = slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME
	title = normalized_text(frames[title_role])
	outline = normalized_text(frames[outline_role])
	expected_title = f"{movie_data.title} ({movie_data.year})"
	require(title == expected_title, "Converted title changed movie identity")
	for label in EXPECTED_LABELS:
		require(label in outline, f"Converted outline is missing literal label: {label}")
	for fragment in expected_outline_fragments(movie_data):
		require(fragment in outline, f"Converted outline is missing movie text: {fragment}")
	lists = list(frames[outline_role].iter(f"{{{TEXT_NS}}}list"))
	list_texts = [normalized_text(text_list) for text_list in lists]
	require(
		any("IMDB rating" in text and "Critics: RT" in text for text in list_texts),
		"Converted outline lost the nested ratings bullet meaning",
	)
	styles = style_map(content_root, styles_root)
	for role in (title_role, outline_role):
		values = referenced_style_values(frames[role], styles)
		require(
			slide_maker.slide_builder.FONT_NAME in values,
			f"Converted role {role} has no OpenDyslexic style reference",
		)
	title_values = referenced_style_values(frames[title_role], styles)
	outline_values = referenced_style_values(frames[outline_role], styles)
	require(
		"shrink-to-fit" in title_values or "true" in title_values,
		"Converted title has no shrink-to-fit autofit property",
	)
	require(
		"shrink-to-fit" in outline_values or "true" in outline_values,
		"Converted outline has no shrink-to-fit autofit property",
	)


#============================================
def frame_geometry_cm(frame: xml.etree.ElementTree.Element) -> dict[str, float]:
	"""Return converted frame geometry in centimeters."""
	geometry = {}
	for odf_name, unused_source_name in GEOMETRY_ATTRIBUTES:
		value = frame.get(f"{{{SVG_NS}}}{odf_name}")
		require(bool(value), f"Converted frame has no {odf_name} geometry")
		geometry[odf_name] = slide_maker.slide_convert.length_cm(value)
	return geometry


#============================================
def validate_converted_geometry(
	page: xml.etree.ElementTree.Element,
	source_geometry: dict[str, dict[str, float]],
	poster_anchor: dict[str, float],
	movie_data: slide_maker.moviedata.MovieData,
) -> None:
	"""Validate role tolerance and converted poster placement against intent."""
	frames = slide_maker.slide_convert.named_frames(page)
	for role in slide_maker.slide_convert.REQUIRED_ROLES:
		converted = frame_geometry_cm(frames[role])
		for geometry_name, unused_source_name in GEOMETRY_ATTRIBUTES:
			difference = abs(converted[geometry_name] - source_geometry[role][geometry_name])
			require(
				difference <= GEOMETRY_TOLERANCE_CM,
				f"Converted role {role} {geometry_name} differs by {difference:.3f} cm",
			)
	poster_role = slide_maker.slide_builder.TEMPLATE_POSTER_NAME
	poster = frame_geometry_cm(frames[poster_role])
	tolerance = GEOMETRY_TOLERANCE_CM
	require(poster["x"] >= poster_anchor["x"] - tolerance, "Converted poster escapes left")
	require(poster["y"] >= poster_anchor["y"] - tolerance, "Converted poster escapes top")
	require(
		poster["x"] + poster["width"] <= poster_anchor["x"] + poster_anchor["width"] + tolerance,
		"Converted poster escapes right",
	)
	require(
		poster["y"] + poster["height"] <= poster_anchor["y"] + poster_anchor["height"] + tolerance,
		"Converted poster escapes bottom",
	)
	anchor_center_x = poster_anchor["x"] + poster_anchor["width"] / 2
	anchor_center_y = poster_anchor["y"] + poster_anchor["height"] / 2
	poster_center_x = poster["x"] + poster["width"] / 2
	poster_center_y = poster["y"] + poster["height"] / 2
	require(abs(anchor_center_x - poster_center_x) <= tolerance, "Converted poster is off-center")
	require(abs(anchor_center_y - poster_center_y) <= tolerance, "Converted poster is off-center")
	with PIL.Image.open(movie_data.poster_path) as poster_image:
		image_ratio = poster_image.width / poster_image.height
	frame_ratio = poster["width"] / poster["height"]
	require(abs(frame_ratio - image_ratio) < 0.01, "Converted poster aspect ratio changed")


#============================================
def run_command(command: list[str], failure_message: str) -> subprocess.CompletedProcess[str]:
	"""Run one render command and preserve its useful diagnostic on failure."""
	result = subprocess.run(command, capture_output=True, text=True, check=False)
	diagnostic = result.stderr.strip() or result.stdout.strip() or "no command diagnostic"
	require(result.returncode == 0, f"{failure_message}: {diagnostic}")
	return result


#============================================
def render_document(
	document_path: pathlib.Path,
	render_path: pathlib.Path,
	work_dir: pathlib.Path,
	tail_fragments: tuple[str, ...],
) -> pathlib.Path:
	"""Render one presentation through PDF and return its accepted PNG path."""
	render_dir = work_dir / f"{document_path.stem}_{document_path.suffix[1:]}_render"
	render_dir.mkdir()
	result = slide_maker.libreoffice_runner.convert_document(
		document_path,
		"pdf",
		render_dir,
		render_dir / "libreoffice_profile",
	)
	diagnostic = result.stderr.strip() or result.stdout.strip() or "no command diagnostic"
	require(result.returncode == 0, f"LibreOffice PDF render failed: {diagnostic}")
	pdf_path = render_dir / f"{document_path.stem}.pdf"
	require(pdf_path.is_file(), f"LibreOffice did not render PDF: {pdf_path}")
	require(pdf_path.stat().st_size > 0, f"LibreOffice rendered an empty PDF: {pdf_path}")
	if tail_fragments:
		text_path = render_dir / f"{document_path.stem}.txt"
		run_command(
			["pdftotext", "-layout", str(pdf_path), str(text_path)],
			"Rendered PDF text extraction failed",
		)
		rendered_text = " ".join(text_path.read_text(encoding="utf-8").split())
		for tail_fragment in tail_fragments:
			require(
				tail_fragment in rendered_text,
				f"Rendered PDF lost end-of-field text: {tail_fragment}",
			)
	render_path.unlink(missing_ok=True)
	run_command(
		[
			"pdftoppm",
			"-png",
			"-singlefile",
			"-r",
			"150",
			str(pdf_path),
			str(render_path.with_suffix("")),
		],
		"Rendered PDF pixel conversion failed",
	)
	require(render_path.is_file(), f"Rendered slide image is absent: {render_path}")
	require(render_path.stat().st_size > 0, f"Rendered slide image is empty: {render_path}")
	with PIL.Image.open(render_path) as render_image:
		render_image.load()
		require(render_image.width > render_image.height, "Rendered slide image is not landscape")
		render_ratio = render_image.width / render_image.height
		expected_ratio = EXPECTED_PAGE_WIDTH_CM / EXPECTED_PAGE_HEIGHT_CM
		require(abs(render_ratio - expected_ratio) < 0.02, "Rendered slide aspect changed")
		statistics = PIL.ImageStat.Stat(render_image.convert("RGB"))
		require(max(statistics.var) > 1.0, "Rendered slide pixels are effectively blank")
	return render_path


#============================================
def convert_control_to_odp(
	control_pptx_path: pathlib.Path,
	work_dir: pathlib.Path,
) -> pathlib.Path:
	"""Convert one empty-text control to temporary ODP without product validation."""
	conversion_dir = work_dir / f"{control_pptx_path.stem}_conversion"
	conversion_dir.mkdir()
	result = slide_maker.libreoffice_runner.convert_document(
		control_pptx_path,
		"odp",
		conversion_dir,
		conversion_dir / "libreoffice_profile",
	)
	diagnostic = result.stderr.strip() or result.stdout.strip() or "no command diagnostic"
	require(
		result.returncode == 0,
		f"LibreOffice empty-text control conversion failed: {diagnostic}",
	)
	control_odp_path = conversion_dir / f"{control_pptx_path.stem}.odp"
	require(control_odp_path.is_file(), "LibreOffice did not convert the empty-text control")
	return control_odp_path


#============================================
def validate_changed_pixel_containment(
	render_path: pathlib.Path,
	control_render_path: pathlib.Path,
	source_geometry: dict[str, dict[str, float]],
	role_name: str,
	role_label: str,
) -> None:
	"""Require pixels changed by one text role to stay inside that role's frame."""
	with PIL.Image.open(render_path) as accepted_source:
		accepted = accepted_source.convert("RGB")
	with PIL.Image.open(control_render_path) as control_source:
		control = control_source.convert("RGB")
	require(accepted.size == control.size, "Accepted and control renders have different dimensions")
	difference = PIL.ImageChops.difference(accepted, control)
	changed_bounds = difference.getbbox()
	require(changed_bounds is not None, f"Accepted render has no {role_label} pixels")
	role_geometry = source_geometry[role_name]
	tolerance = GEOMETRY_TOLERANCE_CM
	left_cm = role_geometry["x"] - tolerance
	top_cm = role_geometry["y"] - tolerance
	right_cm = role_geometry["x"] + role_geometry["width"] + tolerance
	bottom_cm = role_geometry["y"] + role_geometry["height"] + tolerance
	left_px = left_cm * accepted.width / EXPECTED_PAGE_WIDTH_CM
	top_px = top_cm * accepted.height / EXPECTED_PAGE_HEIGHT_CM
	right_px = right_cm * accepted.width / EXPECTED_PAGE_WIDTH_CM
	bottom_px = bottom_cm * accepted.height / EXPECTED_PAGE_HEIGHT_CM
	require(changed_bounds[0] >= left_px, f"Rendered {role_label} escapes left of its frame")
	require(changed_bounds[1] >= top_px, f"Rendered {role_label} escapes above its frame")
	require(changed_bounds[2] <= right_px, f"Rendered {role_label} escapes right of its frame")
	require(changed_bounds[3] <= bottom_px, f"Rendered {role_label} escapes below its frame")
	role_left = int(role_geometry["x"] * accepted.width / EXPECTED_PAGE_WIDTH_CM)
	role_top = int(role_geometry["y"] * accepted.height / EXPECTED_PAGE_HEIGHT_CM)
	role_right = int(
		(role_geometry["x"] + role_geometry["width"])
		* accepted.width
		/ EXPECTED_PAGE_WIDTH_CM
	)
	role_bottom = int(
		(role_geometry["y"] + role_geometry["height"])
		* accepted.height
		/ EXPECTED_PAGE_HEIGHT_CM
	)
	role_box = (role_left, role_top, role_right, role_bottom)
	role_difference = difference.crop(role_box)
	require(role_difference.getbbox() is not None, f"Rendered {role_label} is absent from its frame")


#============================================
def validate_converted_case(
	odp_path: pathlib.Path,
	render_path: pathlib.Path,
	outline_control_odp_path: pathlib.Path,
	title_control_odp_path: pathlib.Path,
	movie_data: slide_maker.moviedata.MovieData,
	source_geometry: dict[str, dict[str, float]],
	poster_anchor: dict[str, float],
	work_dir: pathlib.Path,
) -> None:
	"""Parse, semantically accept, geometrically accept, and render one ODP."""
	archive, content_root, styles_root = parse_odp(odp_path)
	with archive:
		page = validate_page(content_root, styles_root)
		validate_converted_text(page, content_root, styles_root, movie_data)
		slide_maker.slide_convert.validate_poster(archive, page)
		validate_converted_geometry(page, source_geometry, poster_anchor, movie_data)
	tail_fragments = (movie_data.plot[-35:], movie_data.rt_consensus[-35:])
	render_document(odp_path, render_path, work_dir, tail_fragments)
	outline_control_render_path = work_dir / f"{odp_path.stem}_empty_outline.png"
	title_control_render_path = work_dir / f"{odp_path.stem}_empty_title.png"
	render_document(outline_control_odp_path, outline_control_render_path, work_dir, ())
	render_document(title_control_odp_path, title_control_render_path, work_dir, ())
	validate_changed_pixel_containment(
		render_path,
		outline_control_render_path,
		source_geometry,
		slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME,
		"outline text",
	)
	validate_changed_pixel_containment(
		render_path,
		title_control_render_path,
		source_geometry,
		slide_maker.slide_builder.TEMPLATE_TITLE_NAME,
		"title text",
	)


#============================================
def main() -> None:
	"""Build, convert, parse, render, and accept normal and long-text cases."""
	template_path = REPO_ROOT / "template" / "movie_slide_template.pptx"
	require(template_path.is_file(), f"Movie slide template is absent: {template_path}")
	EVIDENCE_DIR.mkdir(parents=True, exist_ok=True)
	with tempfile.TemporaryDirectory() as temporary_directory:
		work_dir = pathlib.Path(temporary_directory)
		case_names = ("her_2013", "her_2013_long_text")
		for case_name, movie_data in zip(case_names, movie_cases(work_dir), strict=True):
			scratch_path = work_dir / f"{case_name}.pptx"
			outline_control_pptx_path = work_dir / f"{case_name}_empty_outline.pptx"
			title_control_pptx_path = work_dir / f"{case_name}_empty_title.pptx"
			odp_path = work_dir / f"{case_name}.odp"
			render_path = EVIDENCE_DIR / f"{case_name}.png"
			source_geometry, poster_anchor = build_and_validate_source(
				movie_data,
				template_path,
				scratch_path,
			)
			write_empty_text_control(
				scratch_path,
				outline_control_pptx_path,
				movie_data,
				slide_maker.slide_builder.TEMPLATE_OUTLINE_NAME,
				"outline",
			)
			write_empty_text_control(
				scratch_path,
				title_control_pptx_path,
				movie_data,
				slide_maker.slide_builder.TEMPLATE_TITLE_NAME,
				"title",
			)
			outline_control_odp_path = convert_control_to_odp(
				outline_control_pptx_path,
				work_dir,
			)
			title_control_odp_path = convert_control_to_odp(
				title_control_pptx_path,
				work_dir,
			)
			slide_maker.slide_convert.convert_presentation(scratch_path, odp_path)
			require(not scratch_path.exists(), "Accepted scratch PPTX was not cleaned")
			validate_converted_case(
				odp_path,
				render_path,
				outline_control_odp_path,
				title_control_odp_path,
				movie_data,
				source_geometry,
				poster_anchor,
				work_dir,
			)
			print(f"Visual acceptance passed for {case_name}: {render_path}")
	print(f"Visual acceptance E2E passed: {EVIDENCE_DIR}")


if __name__ == "__main__":
	main()
